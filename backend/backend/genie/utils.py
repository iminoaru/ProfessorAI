import os
import tempfile
import httpx
import logging
import yt_dlp
import re
from urllib.parse import urlparse
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup
from markdownify import markdownify as md
from pydub import AudioSegment
from openai import OpenAI
from dotenv import load_dotenv
from pytube import extract
import gdown
import mimetypes
from .pdf_utils import process_pdf

# Configure logging
logging.basicConfig(
    level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s"
)
logger = logging.getLogger(__name__)


# Website-specific processors
async def process_youtube(url: str, openai_client: OpenAI) -> str:
    """Process YouTube videos by downloading audio and transcribing"""
    logger.info(f"Processing YouTube URL: {url}")

    try:
        video_id = extract.video_id(url)
        with tempfile.TemporaryDirectory() as temp_dir:
            # Configure yt-dlp options
            ydl_opts = {
                "format": "bestaudio/best",
                "postprocessors": [
                    {
                        "key": "FFmpegExtractAudio",
                        "preferredcodec": "mp3",
                        "preferredquality": "192",
                    }
                ],
                "outtmpl": os.path.join(temp_dir, f"{video_id}.%(ext)s"),
                "quiet": True,
                "no_warnings": True,
            }

            # Download audio
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                ydl.download([url])

            # Process the downloaded audio
            audio_path = os.path.join(temp_dir, f"{video_id}.mp3")
            if not os.path.exists(audio_path):
                raise ValueError("Failed to download audio")

            # Process audio through the audio processor
            with open(audio_path, "rb") as audio_file:
                return process_audio(audio_file.read(), openai_client)

    except Exception as e:
        logger.error(f"Error processing YouTube video: {str(e)}")
        raise ValueError(f"YouTube processing error: {str(e)}")


async def process_arxiv(url: str) -> str:
    """Process arXiv papers by downloading PDF and extracting text"""
    logger.info(f"Processing arXiv URL: {url}")

    try:
        # Extract arXiv ID
        parsed_url = urlparse(url)
        arxiv_id_match = re.search(
            r"/(?:abs|pdf|ps|src|tb)/(?:hep-th/)?(\d+\.\d+|\d+)", parsed_url.path
        )
        if not arxiv_id_match:
            raise ValueError("Invalid arXiv URL format")

        arxiv_id = arxiv_id_match.group(1)
        pdf_url = f"https://arxiv.org/pdf/{arxiv_id}.pdf"

        # Download PDF
        async with httpx.AsyncClient() as client:
            response = await client.get(pdf_url)
            response.raise_for_status()

            return process_pdf(response.content)

    except Exception as e:
        logger.error(f"Error processing arXiv paper: {str(e)}")
        raise ValueError(f"arXiv processing error: {str(e)}")


async def process_github(url: str) -> str:
    """Process GitHub files by converting to raw content and processing"""
    logger.info(f"Processing GitHub URL: {url}")

    try:
        # Convert to raw URL
        raw_url = url.replace("github.com", "raw.githubusercontent.com")
        raw_url = raw_url.replace("/blob/", "/")

        async with httpx.AsyncClient() as client:
            response = await client.get(raw_url)
            response.raise_for_status()
            content = response.content

            # Handle based on file extension
            if url.endswith(".py"):
                return content.decode("utf-8")
            else:
                return md(content.decode("utf-8"))

    except Exception as e:
        logger.error(f"Error processing GitHub file: {str(e)}")
        raise ValueError(f"GitHub processing error: {str(e)}")


# Add this new function after the other specific processors:
async def process_google_drive(url: str, openai_client: OpenAI) -> str:
    """Process Google Drive files by downloading and processing based on type"""
    logger.info(f"Processing Google Drive URL: {url}")

    try:
        with tempfile.TemporaryDirectory() as temp_dir:
            # Download file using gdown
            output_path = os.path.join(temp_dir, "downloaded_file")
            gdown.download(url, output_path, fuzzy=True, quiet=False)

            if not os.path.exists(output_path):
                raise ValueError("Failed to download file from Google Drive")

            # Read file content
            with open(output_path, "rb") as f:
                content = f.read()

            # Detect MIME type
            mime_type, _ = mimetypes.guess_type(output_path)
            if not mime_type:
                # Try to guess from content if no extension
                if content.startswith(b"%PDF"):
                    mime_type = "application/pdf"
                elif content.startswith(b"PK"):
                    # Could be DOCX/PPTX/XLSX, default to DOCX
                    mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                else:
                    mime_type = "text/plain"

            logger.info(f"Detected MIME type for Google Drive file: {mime_type}")
            return process_content(mime_type, content, openai_client)

    except Exception as e:
        logger.error(f"Error processing Google Drive file: {str(e)}")
        raise ValueError(f"Google Drive processing error: {str(e)}")


# Content type processors
def process_html(content: bytes) -> str:
    """Extract text from HTML content"""
    soup = BeautifulSoup(content, "html.parser")
    markdown = md(soup.prettify(), strip=["meta", "link"], newline_style="backslash")
    return re.sub(r"(\n\s*)+\n", "\n\n", markdown)


def process_docx(content: bytes) -> str:
    """Extract text from DOCX content"""
    with tempfile.NamedTemporaryFile(suffix=".docx") as temp_file:
        temp_file.write(content)
        temp_file.flush()
        doc = Document(temp_file.name)
        return "\n\n".join(para.text for para in doc.paragraphs)


def process_pptx(content: bytes) -> str:
    """Extract text from PPTX content"""
    with tempfile.NamedTemporaryFile(suffix=".pptx") as temp_file:
        temp_file.write(content)
        temp_file.flush()
        ppt = Presentation(temp_file.name)
        text_content = []
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
        return "\n\n".join(text_content)


def process_audio(content: bytes, openai_client: OpenAI) -> str:
    """Process and transcribe audio content"""
    with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as temp_file:
        temp_file.write(content)
        temp_file_path = temp_file.name

    try:
        # Truncate to 8 minutes
        eight_minutes_ms = 8 * 60 * 1000
        mp3_audio = AudioSegment.from_file(temp_file_path, format="mp3")
        truncated_audio = mp3_audio[:eight_minutes_ms]
        truncated_audio.export(temp_file_path, format="mp3")

        with open(temp_file_path, "rb") as audio_file:
            transcription = openai_client.audio.transcriptions.create(
                model="whisper-1", file=audio_file
            )

        return transcription.text
    finally:
        os.unlink(temp_file_path)  # Clean up temp file


def process_content(content_type: str, content: bytes, openai_client: OpenAI) -> str:
    """Process content based on its MIME type using pattern matching"""
    # Get base MIME type without parameters
    base_type = content_type.split(";")[0].strip()

    try:
        match base_type:
            case "text/html":
                return process_html(content)
            case "application/pdf":
                return process_pdf(content)
            case "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                return process_docx(content)
            case "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                return process_pptx(content)
            case typ if typ.startswith("audio/"):
                return process_audio(content, openai_client)
            case _:
                raise ValueError(f"Unsupported content type: {base_type}")

    except Exception as e:
        logger.error(f"Error processing content of type {base_type}: {str(e)}")
        raise


async def ingest_url(url: str) -> str:
    """
    Main entry point: Process URLs with special handling for specific websites
    and fallback to content-type based processing
    """
    # Initialize OpenAI client
    load_dotenv()
    openai_api_key = os.getenv("OPENAI_API_KEY")
    if not openai_api_key:
        raise ValueError("OPENAI_API_KEY not found in environment variables")

    openai_client = OpenAI(api_key=openai_api_key)

    try:
        # Parse URL for domain matching
        parsed_url = urlparse(url)
        domain = parsed_url.netloc.lower()

        # Pattern match on domain for special handling
        match domain:
            case d if "youtube.com" in d or "youtu.be" in d:
                return await process_youtube(url, openai_client)
            case d if "arxiv.org" in d:
                return await process_arxiv(url)
            case d if "github.com" in d:
                return await process_github(url)
            case d if "drive.google.com" in d:
                return await process_google_drive(url, openai_client)
            case _:
                # Default handling for other URLs
                async with httpx.AsyncClient() as client:
                    response = await client.get(url, follow_redirects=True)
                    response.raise_for_status()

                    content_type = response.headers.get("Content-Type", "text/html")
                    content = response.content

                    logger.info(
                        f"Processing URL: {response.url} with content type: {content_type}"
                    )
                    return process_content(content_type, content, openai_client)

    except httpx.HTTPError as e:
        logger.error(f"HTTP error while fetching {url}: {str(e)}")
        raise ValueError(f"Error fetching URL: {str(e)}")
    except Exception as e:
        logger.error(f"Error processing URL {url}: {str(e)}")
        raise ValueError(f"Error processing URL: {str(e)}")
