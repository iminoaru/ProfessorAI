from fastapi import APIRouter, HTTPException
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from typing import List
from backend.supabase_client import supabase
from pptx import Presentation
import io
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

router = APIRouter()


class Lesson(BaseModel):
    lesson_id: str
    course_id: str
    chunk_id: str
    title: str
    subtitle: str
    bullet_points: List[str]


def generate_lessons_pptx(lessons: List[Lesson]):
    """
    Generates a professionally designed PPTX file that includes all lessons.
    """
    prs = Presentation()
    
    # Define theme colors and styles
    THEME = {
        'primary': {
            'dark': RGBColor(31, 73, 125),    # Dark blue
            'light': RGBColor(68, 114, 196)   # Light blue
        },
        'accent': {
            'orange': RGBColor(247, 150, 70),  # Orange
            'teal': RGBColor(68, 196, 164)    # Teal
        },
        'background': {
            'light': RGBColor(242, 242, 242), # Light gray
            'white': RGBColor(255, 255, 255)  # White
        },
        'text': {
            'dark': RGBColor(51, 51, 51),     # Dark gray
            'light': RGBColor(255, 255, 255)  # White
        }
    }

    # Style definitions
    STYLES = {
        'title': {
            'font_name': 'Calibri Light',
            'font_size': Pt(44),
            'color': THEME['primary']['dark']
        },
        'subtitle': {
            'font_name': 'Calibri Light',
            'font_size': Pt(28),
            'color': THEME['primary']['light']
        },
        'body': {
            'font_name': 'Calibri',
            'font_size': Pt(18),
            'color': THEME['text']['dark']
        }
    }

    def apply_text_style(paragraph, style):
        """Helper function to apply text styling"""
        font = paragraph.font
        font.name = style['font_name']
        font.size = style['font_size']
        font.color.rgb = style['color']

    # Create title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    # Add a modern geometric background shape
    background = title_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(5), Inches(10), Inches(2.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = THEME['primary']['dark']
    background.line.fill.background()

    

    # Style title and subtitle
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    title.text = "Course Lessons"
    apply_text_style(title.text_frame.paragraphs[0], STYLES['title'])
    
    subtitle.text = f"Total Lessons: {len(lessons)}"
    apply_text_style(subtitle.text_frame.paragraphs[0], STYLES['subtitle'])

    for lesson in lessons:
        # Create section title slide
        section_slide = prs.slides.add_slide(prs.slide_layouts[2])
        
        # Add side accent bar
        accent_bar = section_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), Inches(7.5)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = THEME['accent']['teal']
        accent_bar.line.fill.background()

        # Style section title
        title = section_slide.shapes.title
        title.text = lesson.title
        apply_text_style(title.text_frame.paragraphs[0], STYLES['title'])

        # Add subtitle with modern styling
        subtitle_shape = section_slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(1)
        )
        subtitle_frame = subtitle_shape.text_frame
        subtitle_frame.text = lesson.subtitle
        apply_text_style(subtitle_frame.paragraphs[0], STYLES['subtitle'])

        # Create content slides for bullet points
        bullet_points = lesson.bullet_points
        for i in range(0, len(bullet_points), 4):
            chunk = bullet_points[i:i + 4]
            content_slide = prs.slides.add_slide(prs.slide_layouts[1])

            # Add subtle background pattern
            for j in range(5):
                circle = content_slide.shapes.add_shape(
                    MSO_SHAPE.OVAL, Inches(-2 + j), Inches(j), Inches(0.1), Inches(0.1)
                )
                circle.fill.solid()
                circle.fill.fore_color.rgb = THEME['background']['light']
                circle.line.fill.background()

            # Style content
            title = content_slide.shapes.title
            title.text = f"{lesson.title}" if i == 0 else f"{lesson.title} (cont'd)"
            apply_text_style(title.text_frame.paragraphs[0], STYLES['title'])

            body = content_slide.placeholders[1]
            tf = body.text_frame
            tf.clear()

            for point in chunk:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
                apply_text_style(p, STYLES['body'])
                p.space_before = Pt(12)
                p.space_after = Pt(12)

    # Save presentation
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io


def generate_lessons_pdf(lessons: List[Lesson]):
    """
    Generates a PDF file that includes all lessons with an improved PowerPoint-like font theme and style.

    Args:
        lessons (List[Lesson]): List of Lesson objects.

    Returns:
        BytesIO: PDF file in memory.
    """
    # Register Helvetica font
    base_font = "Helvetica"
    bold_font = "Helvetica-Bold"

    # Create a PDF in memory
    pdf_io = io.BytesIO()
    doc = SimpleDocTemplate(
        pdf_io,
        pagesize=letter,
        topMargin=72,
        bottomMargin=72,
        leftMargin=72,
        rightMargin=72,
    )

    # Define custom styles with better font sizes and colors for a PowerPoint-like appearance
    styles = getSampleStyleSheet()

    title_style = ParagraphStyle(
        "Title",
        parent=styles["Title"],
        fontName=bold_font,
        fontSize=26,  # Larger font size for slide titles
        leading=30,  # Line spacing
        spaceAfter=18,  # Space after title
        textColor=colors.HexColor("#003366"),  # Dark blue color for titles
    )

    subtitle_style = ParagraphStyle(
        "Subtitle",
        parent=styles["Heading2"],
        fontName=base_font,
        fontSize=18,  # Slightly larger font size for subtitles
        leading=22,  # Line spacing
        spaceAfter=12,  # Space after subtitle
        textColor=colors.HexColor("#333333"),  # Dark gray color for subtitles
    )

    bullet_point_style = ParagraphStyle(
        "BulletPointStyle",
        parent=styles["BodyText"],
        fontName=base_font,
        fontSize=14,  # Slightly larger font for better readability
        leading=18,  # Line spacing
        bulletIndent=10,
        leftIndent=36,  # Indent to mimic bullet point structure
        spaceAfter=8,  # Space between bullet points
        textColor=colors.black,  # Standard black color for bullets
    )

    # Container for the PDF elements
    elements = []

    # Add a title page with improved style
    elements.append(Paragraph("Course Lessons", title_style))
    elements.append(Paragraph(f"Total Lessons: {len(lessons)}", subtitle_style))
    elements.append(PageBreak())

    # Iterate over each lesson and create content
    for lesson in lessons:
        # Add the lesson title and subtitle
        elements.append(Paragraph(lesson.title, title_style))
        elements.append(Paragraph(lesson.subtitle, subtitle_style))
        elements.append(Spacer(1, 12))

        # Split bullet points into chunks of 5 to avoid overcrowding the "slide"
        bullet_points = lesson.bullet_points
        bullet_point_chunks = [
            bullet_points[i : i + 5] for i in range(0, len(bullet_points), 5)
        ]

        for chunk_index, chunk in enumerate(bullet_point_chunks):
            if chunk_index > 0:
                # Add a continuation title for additional slides of the same lesson
                elements.append(Paragraph(f"{lesson.title} (cont'd)", subtitle_style))
                elements.append(Spacer(1, 12))

            # Add each bullet point
            for point in chunk:
                elements.append(Paragraph(f"• {point}", bullet_point_style))

            # Add a page break unless it's the last chunk of bullet points for this lesson
            if chunk_index < len(bullet_point_chunks) - 1 or lesson != lessons[-1]:
                elements.append(PageBreak())

    # Build and finalize the PDF
    doc.build(elements)

    # Return the PDF in memory
    pdf_io.seek(0)
    return pdf_io


@router.get("/lessons-pptx/{course_id}", status_code=200)
async def get_lessons_pptx(course_id: str, user_id: str):
    # Fetch all lessons for the course
    res = supabase.table("lessons").select("*").eq("course_id", course_id).execute()
    if not res.data:
        raise HTTPException(status_code=404, detail="No lessons found for this course")
    lessons_data = res.data

    # Convert lessons_data into a list of Lesson objects
    lessons = [Lesson(**lesson_data) for lesson_data in lessons_data]

    # Generate the PPTX file that includes all lessons
    pptx_io = generate_lessons_pptx(lessons)

    # Return the PPTX file as a streaming response
    return StreamingResponse(
        pptx_io,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": f"attachment; filename=Lessons_{course_id}.pptx"
        },
    )


@router.get("/lessons-pdf/{course_id}", status_code=200)
async def get_lessons_pdf(course_id: str, user_id: str):
    # Fetch all lessons for the course
    res = supabase.table("lessons").select("*").eq("course_id", course_id).execute()
    if not res.data:
        raise HTTPException(status_code=404, detail="No lessons found for this course")
    lessons_data = res.data

    # Convert lessons_data into a list of Lesson objects
    lessons = [Lesson(**lesson_data) for lesson_data in lessons_data]

    # Generate the PDF file that includes all lessons
    pdf_io = generate_lessons_pdf(lessons)

    # Return the PDF as a streaming response
    return StreamingResponse(
        pdf_io,
        media_type="application/pdf",
        headers={"Content-Disposition": f"inline; filename=Lessons_{course_id}.pdf"},
    )
